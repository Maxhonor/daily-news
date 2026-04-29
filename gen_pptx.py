#!/usr/bin/env python3
"""Generate NHK news briefing PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_RED = RGBColor(0xE0, 0x3E, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet(slide, left, top, width, height, items, font_size=16, color=WHITE, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox

# === Slide 1: Title ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 2.5, 11, 1.5, "🇯🇵 日本今日要闻", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.2, 11, 0.8, "2026年4月28日（火）| NHK 新闻速报", 22, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.6, "🔥 岩手县大槌町山林火灾持续 · 灭火进行中", 18, YELLOW, False, PP_ALIGN.CENTER)

# === Slide 2: 政治 ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "🏛️ 政治", 36, ACCENT_RED, True)
items = [
    "▸ 高市首相与埃及总统电话会谈，就局势降温达成合作",
    "▸ 冲绳尖阁诸岛（钓鱼岛）附近，中国海警局2艘船已驶出领海",
]
add_bullet(slide, 0.5, 1.5, 12, 4, items, 20)

# === Slide 3: 经济 ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "💰 经济", 36, YELLOW, True)
items = [
    "▸ 日银（日本央行）宣布维持利率不变",
    "▸ 指出「物价上涨风险需要注意」",
    "▸ 本田撤回「2040年新车全部EV化」的目标",
    "▸ 纽约原油突破100美元，创本月新高",
    "▸ 东京燃气警告：受中东局势影响，9月起燃气费可能上涨",
]
add_bullet(slide, 0.5, 1.5, 12, 5, items, 20)

# === Slide 4: 国际 ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "🌍 国际", 36, ACCENT_BLUE, True)
items = [
    "▸ 美国就伊朗新提案展开讨论，但持怀疑态度",
    "▸ 谈判或将继续推进",
    "",
    "🇰🇷 韩国",
    "▸ 尹锡悦前总统夫人金建希二审被判4年",
]
add_bullet(slide, 0.5, 1.5, 12, 5, items, 20)

# === Slide 5: 社会 & 体育 ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 5.5, 0.8, "🏭 社会", 36, RGBColor(0xFF, 0x99, 0x44), True)
items_social = [
    "▸ 广岛制药公司工厂发生储罐爆炸",
    "▸ 5名工人受伤",
    "▸ 新防灾气象信息将于5月28日公布",
]
add_bullet(slide, 0.5, 1.5, 5.5, 4, items_social, 18)

add_text_box(slide, 7, 0.3, 5.5, 0.8, "⚾ 体育", 36, RGBColor(0x66, 0xCC, 0x66), True)
items_sports = [
    "▸ 村上宗隆（白袜队）",
    "▸ 打出第12支本垒打",
    "▸ 领跑两联盟",
]
add_bullet(slide, 7, 1.5, 5.5, 4, items_sports, 20)

# === Slide 6: 山火详情 ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "🔥 焦点：岩手县大槌町山林火灾", 36, RGBColor(0xFF, 0x44, 0x44), True)
items = [
    "▸ 已持续多日的大型山林火灾",
    "▸ 连续两天降雨，灭火工作取得进展",
    "▸ 灭火作业仍在持续中",
    "▸ 暂无人员伤亡报告",
]
add_bullet(slide, 0.5, 1.5, 12, 4, items, 22)

# Save
output_path = "/mnt/c/Users/17932/Desktop/NHK新闻简报_20260428.pptx"
prs.save(output_path)
print(f"✅ PPT saved: {output_path}")
print(f"   Size: {os.path.getsize(output_path)} bytes")
print(f"   Slides: {len(prs.slides)}")
